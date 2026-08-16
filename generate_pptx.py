#!/usr/bin/env python3
import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is not installed. Run 'pip install python-pptx' first.")
    sys.exit(1)

def create_presentation():
    # Initialize Presentation
    prs = Presentation()
    
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Color Palette (Premium Blue/Slate Theme)
    BG_DARK = RGBColor(15, 23, 42)      # Dark Slate Background (#0F172A)
    BG_LIGHT = RGBColor(248, 250, 252) # Off-White Background (#F8FAFC)
    TEXT_DARK = RGBColor(15, 23, 42)    # Near Black Text
    TEXT_LIGHT = RGBColor(255, 255, 255)# Pure White Text
    TEXT_MUTED = RGBColor(100, 116, 139)# Slate Gray Text (#64748B)
    ACCENT_BLUE = RGBColor(10, 54, 157) # Deep Cobalt Blue (#0A369D)
    ACCENT_GREEN = RGBColor(16, 185, 129)# Emerald Green (#10B981)
    CARD_BG = RGBColor(255, 255, 255)   # White Card Fill
    CARD_BORDER = RGBColor(226, 232, 240)# Light Gray Border (#E2E8F0)
    
    # Helper: Set Slide Background
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Create Header
    def add_slide_header(slide, title_text, category_text="AI FOR MANAGERS"):
        # Category label (e.g., AI FOR MANAGERS)
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_right = cat_tf.margin_bottom = 0
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK

    # Helper: Add Custom Card
    def add_card(slide, left, top, width, height, title, content_lines, accent_color=ACCENT_BLUE):
        # Draw background shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1.5)
        
        # Text Box Inside Card
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Card Title
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = accent_color
        p_title.space_after = Pt(10)
        
        # Card Content Lines
        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(6)
            
            # Simple bold formatting parser for text prefixes
            if line.startswith("• "):
                p.level = 0
                if ":" in line:
                    parts = line.split(":", 1)
                    p.text = parts[0] + ":"
                    p.font.bold = True
                    
                    run = p.add_run()
                    run.text = parts[1]
                    run.font.bold = False
                    run.font.size = Pt(12)
                    run.font.color.rgb = TEXT_DARK

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # ==========================================
    # SLIDE 1: Title Slide (Dark Background)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_DARK)
    
    # Large Decorative Circle Accent
    circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1.5), Inches(6.5), Inches(6.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(10, 54, 157)
    circle.line.fill.background()
    
    # Title Text Frame
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    # NMIMS Label
    p_cat = tf1.paragraphs[0]
    p_cat.text = "NMIMS : BUSINESS MANAGEMENT SCHOOL"
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(14)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_GREEN
    p_cat.space_after = Pt(20)
    
    # Title
    p_title = tf1.add_paragraph()
    p_title.text = "Project Framework"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(48)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT
    
    # Subtitle
    p_sub = tf1.add_paragraph()
    p_sub.text = "Travel Agent Fraud & Fake Booking Detector\nHospitality & Travel Vertical"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(10)
    
    # Footer Label
    lbl_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(8.0), Inches(0.8))
    lbl_tf = lbl_box.text_frame
    p_lbl = lbl_tf.paragraphs[0]
    p_lbl.text = "AI OPPORTUNITY ASSESSMENT & BLUEPRINT  |  BUSINESS CONTEXT"
    p_lbl.font.name = "Arial"
    p_lbl.font.size = Pt(11)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: Industry Challenges & Selecting the Topic
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BG_LIGHT)
    add_slide_header(slide2, "Hospitality & Travel: Key Business Challenges")
    
    desc_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8))
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    p_desc = desc_tf.paragraphs[0]
    p_desc.text = "The travel sector operates on perishable inventory. We evaluated three major pain points and selected the highest-value problem to solve using AI:"
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = TEXT_DARK
    
    card_w = Inches(3.64)
    card_h = Inches(4.2)
    card_top = Inches(2.4)
    gap = Inches(0.4)
    left_start = Inches(0.8)
    
    add_card(
        slide2, 
        left_start, 
        card_top, 
        card_w, 
        card_h, 
        "1. Inventory Spoilage", 
        [
            "• What happens: Rogue agents block seats/rooms on hold for free, then cancel them last-minute if they can't sell them.",
            "• Business Impact: The hotel or airline is left with empty rooms/seats that could have been sold to real passengers."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide2, 
        left_start + card_w + gap, 
        card_top, 
        card_w, 
        card_h, 
        "2. Payment Fraud (Selected)", 
        [
            "• What happens: Scammers use stolen cards to buy flights through agent portals. The airline pays for the ticket and faces heavy chargeback fees.",
            "• Business Impact: Direct cash losses and risk of being banned by card networks."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide2, 
        left_start + (card_w + gap) * 2, 
        card_top, 
        card_w, 
        card_h, 
        "3. Commission Abuse", 
        [
            "• What happens: Fake travel agents make bookings just to pocket sales commissions, then cancel the bookings later.",
            "• Business Impact: Payouts are wasted on fake transactions, draining marketing budgets."
        ],
        ACCENT_BLUE
    )

    # ==========================================
    # SLIDE 3: The Double-Agent Scam (Stolen Card Mechanism)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, BG_LIGHT)
    add_slide_header(slide3, "Understanding the Threat: The Double-Agent Scam")
    
    desc_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8))
    desc_tf3 = desc_box3.text_frame
    desc_tf3.word_wrap = True
    p_desc3 = desc_tf3.paragraphs[0]
    p_desc3.text = "Stolen card fraud is NOT rare; it is the most common scam in travel because tickets are high-value and digital. Here is how a fake travel agent exploits the system:"
    p_desc3.font.name = "Arial"
    p_desc3.font.size = Pt(14)
    p_desc3.font.color.rgb = TEXT_DARK
    
    # 3 horizontal step cards
    step_w = Inches(3.64)
    step_h = Inches(4.2)
    step_top = Inches(2.4)
    
    add_card(
        slide3,
        left_start,
        step_top,
        step_w,
        step_h,
        "Step 1: Pocket the Cash",
        [
            "• Scenario: A real customer visits a fake travel agent and books a flight for $1,000.",
            "• Action: The customer pays the agent in cash or direct bank transfer.",
            "• Result: The fake agent now has $1,000 cash in hand."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide3,
        left_start + step_w + gap,
        step_top,
        step_w,
        step_h,
        "Step 2: Book with Stolen Card",
        [
            "• Action: The fake agent goes to the dark web, buys a stolen card, and books the real ticket on the airline portal.",
            "• Result: The customer gets a valid ticket. The fake agent has pocketed $1,000 clean profit."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide3,
        left_start + (step_w + gap) * 2,
        step_top,
        step_w,
        step_h,
        "Step 3: The Chargeback Hit",
        [
            "• Action: Weeks later, the real cardholder sees the charge and reports it.",
            "• Penalty: The bank forces the airline to refund the $1,000.",
            "• Damage: The airline loses the ticket value, paid bank fees, and the scammer is gone."
        ],
        ACCENT_GREEN
    )

    # ==========================================
    # SLIDE 4: Research Methodology & Who We Interview
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, BG_LIGHT)
    add_slide_header(slide4, "Primary & Secondary Research Strategy")
    
    sec_w = Inches(5.6)
    sec_h = Inches(5.0)
    
    # Left Column: Secondary
    add_card(
        slide4,
        Inches(0.8),
        Inches(1.6),
        sec_w,
        sec_h,
        "Secondary Research Plan",
        [
            "• Where we look: Industry reports (Stripe, McKinsey), international airline fraud data (IATA), and past tech case studies.",
            "• What we learn: The average cost of chargebacks, how dynamic pricing reacts to fake holds, and how other companies stop it."
        ],
        ACCENT_BLUE
    )
    
    # Right Column: Primary (Clarifying Who We Interview)
    add_card(
        slide4,
        Inches(0.8) + sec_w + Inches(0.533),
        Inches(1.6),
        sec_w,
        sec_h,
        "Primary Research (Whom Do We Interview?)",
        [
            "• NOTE: We DO NOT interview the criminals/fraudsters.",
            "• We interview the victims and managers suffering from fraud:",
            "  - 1. Fraud Operations Managers: The business defenders who lose sleep checking logs and dealing with card chargebacks.",
            "  - 2. Honest Travel Agents: Portal users who suffer when security filters become too slow or lock them out of bookings.",
            "  - 3. End Consumers: Travelers who face sudden ticket cancellations or price hikes.",
            "• Data target: 3+ detailed interviews or 10+ survey responses."
        ],
        ACCENT_GREEN
    )

    # ==========================================
    # SLIDE 5: Opportunity Prioritization Matrix (Simplified Table)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, BG_LIGHT)
    add_slide_header(slide5, "AI Opportunity Evaluation Matrix")
    
    lbl = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
    p_lbl = lbl.text_frame.paragraphs[0]
    p_lbl.text = "Comparing potential AI projects to justify our focus. Fraud Detection offers the highest direct cost savings and best data availability:"
    p_lbl.font.name = "Arial"
    p_lbl.font.size = Pt(13)
    p_lbl.font.color.rgb = TEXT_MUTED

    # Add Table
    rows = 4
    cols = 7
    left = Inches(0.8)
    top = Inches(2.2)
    width = Inches(11.7)
    height = Inches(4.0)
    
    table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Columns Width
    table.columns[0].width = Inches(2.7) # Opportunity
    table.columns[1].width = Inches(1.5) # Impact
    table.columns[2].width = Inches(1.5) # Money Saved
    table.columns[3].width = Inches(1.5) # Sales Potential
    table.columns[4].width = Inches(1.5) # Ease of Building
    table.columns[5].width = Inches(1.5) # Data Available
    table.columns[6].width = Inches(1.5) # Final Score
    
    headers = ["Project Idea", "Business Value", "Money Saved", "Sales Potential", "Ease of Building", "Data Available", "Priority Score"]
    row_data = [
        ["AI Travel Fraud Detector", "High", "High", "Medium", "Medium", "High", "4.2 (Selected)"],
        ["Dynamic Ticket Pricing", "High", "Low", "High", "Hard", "High", "3.8"],
        ["AI Travel Itinerary Planner", "Medium", "Low", "High", "Easy", "Medium", "3.4"]
    ]
    
    # Format Headers
    for c in range(cols):
        cell = table.cell(0, c)
        cell.text = headers[c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        
    # Format Rows
    for r in range(len(row_data)):
        data = row_data[r]
        for c in range(cols):
            cell = table.cell(r + 1, c)
            cell.text = data[c]
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = RGBColor(239, 246, 255) # Highlight selected row
            else:
                cell.fill.fore_color.rgb = CARD_BG
                
            p = cell.text_frame.paragraphs[0]
            if c == 0:
                p.alignment = PP_ALIGN.LEFT
                p.font.bold = True
            else:
                p.alignment = PP_ALIGN.CENTER
            
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            
            if c == 6:
                p.font.bold = True
                if r == 0:
                    p.font.color.rgb = ACCENT_BLUE

    # ==========================================
    # SLIDE 6: AI Solution: Explained via Simple Analogies
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, BG_LIGHT)
    add_slide_header(slide6, "How the AI Works: Our 'Smart Security' Team")
    
    step_w = Inches(2.7)
    step_h = Inches(4.5)
    step_top = Inches(1.8)
    step_gap = Inches(0.3)
    step_start = Inches(0.8)
    
    add_card(
        slide6, 
        step_start, 
        step_top, 
        step_w, 
        step_h, 
        "1. The Gatekeeper", 
        [
            "• Role: Rule Checker.",
            "• Simple Analogy: Checks ID card details at the door.",
            "• What it does: Flags bookings where the card's billing country doesn't match where the user is logging in from."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide6, 
        step_start + (step_w + step_gap), 
        step_top, 
        step_w, 
        step_h, 
        "2. The Shop Assistant", 
        [
            "• Role: Speed Monitor.",
            "• Simple Analogy: Watches how fast someone moves in a shop.",
            "• What it does: Flags bookings typed at superhuman speeds (in milliseconds), identifying automated software bots."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide6, 
        step_start + (step_w + step_gap) * 2, 
        step_top, 
        step_w, 
        step_h, 
        "3. The Detective", 
        [
            "• Role: Connection Finder.",
            "• Simple Analogy: Connects dots on a whiteboard.",
            "• What it does: Detects if 5 different travel agents are sharing the same credit card number or device, exposing fraud networks."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide6, 
        step_start + (step_w + step_gap) * 3, 
        step_top, 
        step_w, 
        step_h, 
        "4. The Auditor & Manager", 
        [
            "• Role: Language Reader.",
            "• Simple Analogy: Reads note logs in files.",
            "• What it does: Scans text comments for bot templates. Then, drafts a simple warning summary for our human employees."
        ],
        ACCENT_GREEN
    )

    # ==========================================
    # SLIDE 7: Implementing the Gatekeeper (Logical Flow)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, BG_LIGHT)
    add_slide_header(slide7, "Implementation: How 'The Gatekeeper' Works")
    
    desc_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8))
    desc_tf7 = desc_box7.text_frame
    desc_tf7.word_wrap = True
    p_desc7 = desc_tf7.paragraphs[0]
    p_desc7.text = "The Gatekeeper runs instantly in the portal backend when a user clicks 'Book Now'. Here is the exact logical step-by-step program:"
    p_desc7.font.name = "Arial"
    p_desc7.font.size = Pt(14)
    p_desc7.font.color.rgb = TEXT_DARK
    
    add_card(
        slide7,
        left_start,
        step_top,
        step_w,
        step_h,
        "Step 1: Extract Data",
        [
            "• IP Address: Checks where the agent is logging in (e.g. Russia).",
            "• Card Details: Looks up card's Bank BIN number country (e.g. India).",
            "• Flight Time: Checks when the flight departs (e.g. in 1.5 hours)."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide7,
        left_start + step_w + gap,
        step_top,
        step_w,
        step_h,
        "Step 2: Run Logic Check",
        [
            "• Test A: Does the IP country match the Card country? (No, Russia vs India).",
            "• Test B: Is flight time departure less than 2 hours? (Yes, 1.5 hours left).",
            "• Result: High mismatch flag."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide7,
        left_start + (step_w + gap) * 2,
        step_top,
        step_w,
        step_h,
        "Step 3: Trigger Action",
        [
            "• Score Calculated: System assigns 95/100 risk score.",
            "• Transaction Blocked: Block ticket issuance instantly.",
            "• Alert: Recommends security verification, protecting the airline."
        ],
        ACCENT_GREEN
    )

    # ==========================================
    # SLIDE 8: Dual-Value Proposition (Who Wins?)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, BG_LIGHT)
    add_slide_header(slide8, "AI Dual-Value Proposition Strategy")
    
    card_w = Inches(5.6)
    card_h = Inches(5.0)
    
    add_card(
        slide8,
        Inches(0.8),
        Inches(1.6),
        card_w,
        card_h,
        "Value for the Travel Company",
        [
            "• Stop Financial Leakage: Prevents expensive credit card chargebacks and bank fines.",
            "• Protect Inventory: Releases seats/rooms held by fake bookings, ensuring they can be sold to real customers.",
            "• Save Staff Labor: Generates easy-to-read warnings for operations staff so they don't waste time on manual checks.",
            "• Target Goals: 70% fewer chargebacks; 15% better room/seat occupancy."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide8,
        Inches(0.8) + card_w + Inches(0.533),
        Inches(1.6),
        card_w,
        card_h,
        "Value for the Everyday Customer",
        [
            "• Cheaper Bookings: Stops bot nets from artificially inflating ticket demand, keeping flight prices fair.",
            "• Scam Prevention: Flags compromised travel agent credentials before they can sell fake boarding tickets to customers.",
            "• Travel Security: Ensures customer tickets are completely verified, preventing gate cancellations at the airport."
        ],
        ACCENT_GREEN
    )

    # ==========================================
    # SLIDE 9: Risks & Ethical Concerns
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, BG_LIGHT)
    add_slide_header(slide9, "Critical Risks & Ethical Concerns")
    
    card_w = Inches(3.64)
    card_h = Inches(4.5)
    card_top = Inches(1.8)
    gap = Inches(0.4)
    left_start = Inches(0.8)
    
    add_card(
        slide9, 
        left_start, 
        card_top, 
        card_w, 
        card_h, 
        "1. Protecting Privacy", 
        [
            "• The Risk: Storing passenger names, card data, and click history raises severe privacy concerns.",
            "• The Fix: Hashing/encrypting personal details before the AI analyzes them. Storing payment data only in secure vaults."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide9, 
        left_start + card_w + gap, 
        card_top, 
        card_w, 
        card_h, 
        "2. Geographical Bias", 
        [
            "• The Risk: The AI could block honest travel agents in developing countries because their cards or booking habits look unusual.",
            "• The Fix: Adjusting the AI settings to judge risk based on local behaviors, not one-size-fits-all rules."
        ],
        ACCENT_BLUE
    )
    
    add_card(
        slide9, 
        left_start + (card_w + gap) * 2, 
        card_top, 
        card_w, 
        card_h, 
        "3. Clear Explanations", 
        [
            "• The Risk: If the AI is a 'black box', human staff won't know *why* a booking was flagged and might make mistakes.",
            "• The Fix: The AI Manager translates its logic into a simple sentence (e.g., 'Flagged for inhuman booking speed')."
        ],
        ACCENT_BLUE
    )

    # ==========================================
    # SLIDE 10: Step-by-Step Project Roadmap (Planning to Execution)
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, BG_DARK)
    
    # Large Decorative Circle Accent
    circle = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(3.5), Inches(6.5), Inches(6.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(10, 54, 157)
    circle.line.fill.background()
    
    # Title Text Frame
    text_box = slide10.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5))
    tf10 = text_box.text_frame
    tf10.word_wrap = True
    
    p10_title = tf10.paragraphs[0]
    p10_title.text = "Planning & Execution Timeline"
    p10_title.font.name = "Arial"
    p10_title.font.size = Pt(36)
    p10_title.font.bold = True
    p10_title.font.color.rgb = TEXT_LIGHT
    p10_title.space_after = Pt(20)
    
    p10_steps = [
        "1. Planning & Data Gathering (Months 1-2): Collect past booking logs and label fake holds vs. real transactions.",
        "2. Building the AI Security (Months 3-4): Code the logic for the Gatekeeper, Speed Monitor, and Detective models.",
        "3. Live Execution & Launch (Month 5): Connect the AI models to run in <150ms at portal checkout.",
        "4. Human Integration (Month 6): Deploy the GenAI Co-pilot to summarize flags and draft agent emails.",
        "5. Continuous Learning (Ongoing): Train the AI with human-verified logs so it remains smart."
    ]
    
    for step in p10_steps:
        p = tf10.add_paragraph()
        p.text = step
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)
        
        # Bold the prefix step number
        parts = step.split(":", 1)
        p.text = parts[0] + ":"
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        
        run = p.add_run()
        run.text = parts[1]
        run.font.bold = False
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT_LIGHT
    
    # Save Presentation
    output_filename = "Travel_Fraud_AI_Blueprint_Presentation.pptx"
    output_path = os.path.join(os.path.dirname(__file__), output_filename)
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
